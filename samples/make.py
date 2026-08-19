import os, zipfile, docx, openpyxl
from openpyxl import Workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
os.chdir("samples")

# ---- DOCX ----
d = docx.Document()
d.add_heading("Quarterly Review", 0)
d.add_heading("Revenue", 1)
d.add_paragraph("Revenue grew fourteen percent across every region this year.")
d.add_paragraph("Northern territory outperformed expectations.", style="List Bullet")
d.add_paragraph("Southern territory met target.", style="List Bullet")
d.add_heading("Risks", 1)
d.add_paragraph("Supply chain delays remain the largest open risk.")
t = d.add_table(rows=3, cols=3)
for i,v in enumerate(["Region","Owner","Status"]): t.rows[0].cells[i].text = v
for i,v in enumerate(["North","Priya","Green"]): t.rows[1].cells[i].text = v
for i,v in enumerate(["South","Arun","Amber"]): t.rows[2].cells[i].text = v
d.save("sample.docx")

# ---- XLSX ----
wb = Workbook(); ws = wb.active; ws.title = "Q1"
for r in [["Item","Qty","Price"],["Widget",10,2.5],["Gadget",4,19.0],["Sprocket",7,3.25]]: ws.append(r)
ws2 = wb.create_sheet("Q2")
for r in [["Item","Qty"],["Widget",12]]: ws2.append(r)
wb.save("sample.xlsx")

# ---- PDF ----
c = canvas.Canvas("sample.pdf", pagesize=A4)
c.setFont("Helvetica-Bold", 20); c.drawString(72, 780, "Annual Summary")
c.setFont("Helvetica", 11)
c.drawString(72, 740, "Revenue rose across every region this year.")
c.drawString(72, 722, "Headcount stayed flat at 240 people.")
c.setFont("Helvetica-Bold", 14); c.drawString(72, 690, "Outlook")
c.setFont("Helvetica", 11); c.drawString(72, 668, "Growth is expected to continue into next year.")
c.showPage()
c.setFont("Helvetica", 11); c.drawString(72, 780, "Second page content here.")
c.showPage(); c.save()

# ---- PPTX ----
try:
    from pptx import Presentation
    from pptx.util import Inches
    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[1])
    s.shapes.title.text = "Project Alpha"
    s.placeholders[1].text = "Kickoff meeting\nThree regions in scope"
    s.notes_slide.notes_text_frame.text = "Remember to mention the budget."
    s2 = p.slides.add_slide(p.slide_layouts[1])
    s2.shapes.title.text = "Next Steps"
    s2.placeholders[1].text = "Hire two engineers\nShip by December"
    p.save("sample.pptx")
    print("pptx ok")
except Exception as e:
    print("pptx skipped:", e)

# ---- ODT (built by hand, no library needed) ----
content = """<?xml version="1.0" encoding="UTF-8"?>
<office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0" xmlns:table="urn:oasis:names:tc:opendocument:xmlns:table:1.0">
<office:body><office:text>
<text:h text:outline-level="1">Meeting Notes</text:h>
<text:p>Attendees agreed the schedule is realistic.</text:p>
<text:h text:outline-level="2">Actions</text:h>
<text:list><text:list-item><text:p>Draft the proposal</text:p></text:list-item><text:list-item><text:p>Book the venue</text:p></text:list-item></text:list>
<table:table><table:table-row><table:table-cell><text:p>Task</text:p></table:table-cell><table:table-cell><text:p>Owner</text:p></table:table-cell></table:table-row><table:table-row><table:table-cell><text:p>Proposal</text:p></table:table-cell><table:table-cell><text:p>Meera</text:p></table:table-cell></table:table-row></table:table>
</office:text></office:body></office:document-content>"""
with zipfile.ZipFile("sample.odt","w") as z:
    z.writestr("mimetype","application/vnd.oasis.opendocument.text")
    z.writestr("content.xml", content)

# ---- EPUB (built by hand) ----
with zipfile.ZipFile("sample.epub","w") as z:
    z.writestr("mimetype","application/epub+zip")
    z.writestr("META-INF/container.xml",'<?xml version="1.0"?><container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container"><rootfiles><rootfile full-path="OEBPS/book.opf" media-type="application/oebps-package+xml"/></rootfiles></container>')
    z.writestr("OEBPS/book.opf",'<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="id"><metadata xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>The Test Book</dc:title><dc:creator>A. Writer</dc:creator></metadata><manifest><item id="c1" href="ch1.xhtml" media-type="application/xhtml+xml"/><item id="c2" href="ch2.xhtml" media-type="application/xhtml+xml"/></manifest><spine><itemref idref="c1"/><itemref idref="c2"/></spine></package>')
    z.writestr("OEBPS/ch1.xhtml",'<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><body><h1>Chapter One</h1><p>It was a <strong>bright</strong> cold day in April.</p></body></html>')
    z.writestr("OEBPS/ch2.xhtml",'<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><body><h1>Chapter Two</h1><p>The second chapter begins here.</p></body></html>')

# ---- indented TXT (the original annotation bug) ----
open("sample.txt","w",encoding="utf-8").write("QUARTERLY REPORT\n\n    Revenue grew by 14 percent.\n    The north outperformed.\n\n\tTab indented note.\n")
open("sample.csv","w",encoding="utf-8").write('Name,Qty,Notes\nWidget,10,"has, a comma"\nGadget,4,plain\n')
print("done:", sorted(os.listdir(".")))
